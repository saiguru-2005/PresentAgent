from pptx import Presentation
import os

repo_root = os.path.dirname(os.path.abspath(__file__))
template_path = os.path.join(repo_root, "resource", "templates", "default_template.pptx")

print(f"Smart Cleaning: {template_path}")

if not os.path.exists(template_path):
    print("Template not found!")
    exit(1)

prs = Presentation(template_path)

replacements = {
    "Odissi": "Topic",
    "Tourism": "Domain",
    "Dance": "Subject",
    "Culture": "Categorization",
    "Indian": "Global",
    "Orissa": "Location",
    "Bhubaneswar": "City"
}

def clean_text(text):
    for wrong, right in replacements.items():
        text = text.replace(wrong, right)
        text = text.replace(wrong.upper(), right.upper())
        text = text.replace(wrong.lower(), right.lower())
    return text

count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            # We must iterate paragraphs and runs to preserve formatting
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    original = run.text
                    cleaned = clean_text(original)
                    if original != cleaned:
                        run.text = cleaned
                        count += 1
                        
print(f"Replaced {count} instances of contaminated terms.")
prs.save(template_path)
print(f"Saved smart-cleaned template to {template_path}")
