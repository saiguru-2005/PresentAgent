from pptx import Presentation
import os

repo_root = os.path.dirname(os.path.abspath(__file__))
template_path = os.path.join(repo_root, "resource", "templates", "default_template.pptx")

print(f"Inspecting: {template_path}")

if not os.path.exists(template_path):
    print("Template not found!")
    exit(1)

prs = Presentation(template_path)
found_odissi = False
found_tourism = False

for i, slide in enumerate(prs.slides):
    text_content = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_content.append(shape.text)
    
    full_text = " ".join(text_content).lower()
    if "odissi" in full_text:
        found_odissi = True
        print(f"Slide {i+1}: Found 'Odissi'")
    if "tourism" in full_text:
        found_tourism = True
        print(f"Slide {i+1}: Found 'Tourism'")

if found_odissi or found_tourism:
    print("\nCONFIRMED: The default template contains the unwanted content.")
else:
    print("\nCLEAN: The default template does not contain the unwanted content.")
