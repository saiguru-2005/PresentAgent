from pptx import Presentation
import os
import sys

# Path to the template
repo_root = os.path.dirname(os.path.abspath(__file__))
template_path = os.path.join(repo_root, "resource", "templates", "default_template.pptx")

print(f"Verifying template: {template_path}")

try:
    if not os.path.exists(template_path):
        print("ERROR: Template file does not exist.")
        sys.exit(1)

    prs = Presentation(template_path)
    print(f"SUCCESS: Template loaded. Slide count: {len(prs.slides)}")
    
    # Try adding a slide to ensure layouts work
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    print("SUCCESS: Added a test slide.")
    
except Exception as e:
    print(f"CRITICAL FAILURE: Template is corrupt or unusable. Error: {e}")
    sys.exit(1)
