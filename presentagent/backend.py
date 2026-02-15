import asyncio
print("Initializing PresentAgent Backend...")
print("Loading libraries (this may take 10-20 seconds)...")
import functools
import hashlib
import importlib
import json
import os
import hashlib
import json
import logging
import sys


# ----------------------------------------------------------------------
# ----------------------------------------------------------------------
# LOCAL OLLAMA CONFIGURATION (User Request: No API Keys)
# ----------------------------------------------------------------------
os.environ["OPENAI_API_KEY"] = "ollama"  # Placeholder
os.environ["API_BASE"] = "http://localhost:11434/v1"
os.environ["LANGUAGE_MODEL"] = "qwen2.5:7b"
os.environ["VISION_MODEL"] = "moondream"
os.environ["TEXT_MODEL"] = "nomic-embed-text"
# ----------------------------------------------------------------------
# [FIX] Add project root to sys.path so we can import 'pptagent'
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import traceback
import uuid
import requests
from bs4 import BeautifulSoup
from pptagent.evaluation import Evaluator
import subprocess
import tempfile
from contextlib import asynccontextmanager
from copy import deepcopy
from datetime import datetime
from typing import Optional
import shutil

def ensure_default_template():
    target_dir = pjoin(RUNS_DIR, "pptx", "default_template")
    target_file = pjoin(target_dir, "source.pptx")
    
    # Path to the template in the repo (relative to backend.py which is in /presentagent)
    # Repo root is one level up
    repo_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    source_file = os.path.join(repo_root, "resource", "templates", "default_template.pptx")
    
    if not os.path.exists(target_file):
        try:
            os.makedirs(target_dir, exist_ok=True)
            if os.path.exists(source_file):
                logger.info(f"Copying default template to {target_file}")
                shutil.copy(source_file, target_file)
            else:
                 logger.error(f"Default template source not found at {source_file}")
        except Exception as e:
            logger.error(f"Failed to copy default template: {e}")




# from pdf2image import convert_from_path
from pptx import Presentation as PptxPresentation

from fastapi import (
    FastAPI,
    File,
    Form,
    HTTPException,
    Request,
    UploadFile,
    WebSocket,
    WebSocketDisconnect,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
sys.path.append("../")
import pptagent.induct as induct
import pptagent.pptgen as pptgen
from pptagent.document import Document
from pptagent.model_utils import ModelManager, parse_pdf
from pptagent.multimodal import ImageLabler
from pptagent.presentation import Presentation
from pptagent.utils import Config, get_logger, package_join, pjoin, ppt_to_images_async


async def run_blocking(func, *args, **kw):
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, functools.partial(func, *args, **kw))

async def run_cmd(cmd: list[str]):
    proc = await asyncio.create_subprocess_exec(
        *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
    )
    stdout, stderr = await proc.communicate()
    if proc.returncode != 0:
        raise RuntimeError(f"{' '.join(cmd)}\n{stderr.decode()}")
    return stdout
# ----------------------------------------------------------------------

# constants
DEBUG = True if len(sys.argv) == 1 else False
RUNS_DIR = package_join("runs")
STAGES = ["PPT Parsing", "PDF Parsing", "PPT Analysis", "PPT Generation", "Success!"]


ppt_video_progress_store: dict[str, dict] = {}

models = ModelManager()

@asynccontextmanager
async def lifespan(_: FastAPI):
    assert await models.test_connections(), "Model connection test failed"
    yield

# server
logger = get_logger(__name__)
app = FastAPI(lifespan=lifespan)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
progress_store: dict[str, dict] = {}
active_connections: dict[str, WebSocket] = {}


ppt_video_active_connections: dict[str, WebSocket] = {}


class ProgressManager:
    def __init__(self, task_id: str, stages: list[str], debug: bool = True):
        self.task_id = task_id
        self.stages = stages
        self.debug = debug
        self.failed = False
        self.current_stage = 0
        self.total_stages = len(stages)

    async def report_progress(self):
        assert self.task_id in active_connections, "WebSocket connection closed"
        self.current_stage += 1
        progress = int((self.current_stage / self.total_stages) * 100)
        await send_progress(
            active_connections[self.task_id],
            f"Stage: {self.stages[self.current_stage - 1]}",
            progress,
        )

    async def fail_stage(self, error_message: str):
        await send_progress(
            active_connections[self.task_id],
            f"{self.stages[self.current_stage]} Error: {error_message}",
            100,
        )
        self.failed = True
        active_connections.pop(self.task_id, None)
        if self.debug:
            logger.error(
                f"{self.task_id}: {self.stages[self.current_stage]} Error: {error_message}"
            )

@app.post("/api/upload")
async def create_task(
        pptxFile: UploadFile = File(None),
        pdfFile: UploadFile = File(None),
        topic: str = Form(None),
        url: str = Form(None),
        numberOfPages: int = Form(...),
):
    task_id = datetime.now().strftime("20%y-%m-%d") + "/" + str(uuid.uuid4())
    logger.info(f"task created: {task_id}")
    os.makedirs(pjoin(RUNS_DIR, task_id))
    task = {
        "numberOfPages": numberOfPages,
        "pptx": "default_template",
    }
    # Ensure default template exists if we might use it
    ensure_default_template()

    if pptxFile is not None:
        pptx_blob = await pptxFile.read()
        pptx_md5 = hashlib.md5(pptx_blob).hexdigest()
        task["pptx"] = pptx_md5
        pptx_dir = pjoin(RUNS_DIR, "pptx", pptx_md5)
        if not os.path.exists(pptx_dir):
            os.makedirs(pptx_dir, exist_ok=True)
            with open(pjoin(pptx_dir, "source.pptx"), "wb") as f:
                f.write(pptx_blob)
    if pdfFile is not None:
        pdf_blob = await pdfFile.read()
        pdf_md5 = hashlib.md5(pdf_blob).hexdigest()
        task["pdf"] = pdf_md5
        pdf_dir = pjoin(RUNS_DIR, "pdf", pdf_md5)
        if not os.path.exists(pdf_dir):
            os.makedirs(pdf_dir, exist_ok=True)
            with open(pjoin(pdf_dir, "source.pdf"), "wb") as f:
                f.write(pdf_blob)
    elif url is not None:
        try:
            response = requests.get(url)
            soup = BeautifulSoup(response.text, 'html.parser')
            for script in soup(["script", "style"]):
                script.extract()
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\\n'.join(chunk for chunk in chunks if chunk)
            url_md5 = hashlib.md5(url.encode()).hexdigest()
            task["pdf"] = url_md5
            pdf_dir = pjoin(RUNS_DIR, "pdf", url_md5)
            if not os.path.exists(pdf_dir):
                os.makedirs(pdf_dir, exist_ok=True)
                with open(pjoin(pdf_dir, "source.md"), "w", encoding='utf-8') as f:
                    f.write(text)
        except Exception as e:
            logger.error(f"Failed to parse URL: {e}")
            raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {str(e)}")
    if topic is not None:
        task["pdf"] = topic
    progress_store[task_id] = task
    # Start the PPT generation task asynchronously
    asyncio.create_task(ppt_gen(task_id))
    return {"task_id": task_id.replace("/", "|")}


async def send_progress(websocket: Optional[WebSocket], status: str, progress: int):
    if websocket is None:
        logger.info(f"websocket is None, status: {status}, progress: {progress}")
        return
    await websocket.send_json({"progress": progress, "status": status})


async def send_ppt_video_progress(task_id: str):
    if task_id in ppt_video_active_connections:
        progress_data = ppt_video_progress_store[task_id]
        await ppt_video_active_connections[task_id].send_json(progress_data)
    else:
        logger.warning(f"No WebSocket connection for PPT video task {task_id}")


@app.websocket("/wsapi/{task_id}")
async def websocket_endpoint(websocket: WebSocket, task_id: str):
    task_id = task_id.replace("|", "/")
    if task_id in progress_store:
        await websocket.accept()
    else:
        raise HTTPException(status_code=404, detail="Task not found")
    active_connections[task_id] = websocket
    try:
        while True:
            await websocket.receive_text()
    except WebSocketDisconnect:
        logger.info("websocket disconnected: %s", task_id)
        active_connections.pop(task_id, None)


@app.get("/api/download")
async def download(task_id: str):
    task_id = task_id.replace("|", "/")
    if not os.path.exists(pjoin(RUNS_DIR, task_id)):
        raise HTTPException(status_code=404, detail="Task not created yet")
    file_path = pjoin(RUNS_DIR, task_id, "final.pptx")
    if os.path.exists(file_path):
        return FileResponse(
            file_path,
            media_type="application/pptx",
            headers={"Content-Disposition": "attachment; filename=pptagent.pptx"},
        )
    raise HTTPException(status_code=404, detail="Task not finished yet")


@app.post("/api/feedback")
async def feedback(request: Request):
    body = await request.json()
    feedback = body.get("feedback")
    task_id = body.get("task_id")

    with open(pjoin(RUNS_DIR, "feedback", f"{task_id}.txt"), "w") as f:
        f.write(feedback)
    return {"message": "Feedback submitted successfully"}


@app.get("/")
async def hello():
    return {"message": "Hello, World!"}


@app.post("/api/ppt-to-video")
async def create_ppt_video_task(pptFile: UploadFile = File(...)):
    task_id = str(uuid.uuid4())
    logger.info(f"PPT2Presentation task created: {task_id}")

    task_dir = pjoin(RUNS_DIR, "ppt_video", task_id)
    os.makedirs(task_dir, exist_ok=True)

    ppt_blob = await pptFile.read()
    ppt_path = pjoin(task_dir, "source.pptx")
    with open(ppt_path, "wb") as f:
        f.write(ppt_blob)

    ppt_video_progress_store[task_id] = {
        "status": "processing",
        "current_step": 1,
        "current_slide": 0,
        "total_slides": 0,
        "progress_percentage": 0,
        "task_dir": task_dir,
        "ppt_path": ppt_path
    }

    asyncio.create_task(process_ppt_to_video(task_id))
    return {"task_id": task_id}


@app.websocket("/wsapi/ppt-to-video/{task_id}")
async def websocket_ppt_video_endpoint(websocket: WebSocket, task_id: str):
    if task_id in ppt_video_progress_store:
        await websocket.accept()
    else:
        raise HTTPException(status_code=404, detail="Task not found")
    ppt_video_active_connections[task_id] = websocket
    try:
        while True:
            await websocket.receive_text()
    except WebSocketDisconnect:
        logger.info("PPT video websocket disconnected: %s", task_id)
        ppt_video_active_connections.pop(task_id, None)


async def process_ppt_to_video(task_id: str):
    task_dir = ppt_video_progress_store[task_id]["task_dir"]
    try:
        ppt_path = ppt_video_progress_store[task_id]["ppt_path"]
        ppt_video_progress_store[task_id].update(current_step=1, progress_percentage=10.00)
        await send_ppt_video_progress(task_id)

        pdf_path = pjoin(task_dir, "source.pdf")
        # [MODIFIED] Use absolute path for Windows LibreOffice
        soffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        if not os.path.exists(soffice_path):
             soffice_path = "libreoffice" # Fallback to PATH

        await run_cmd([
            soffice_path, "--headless", "--convert-to", "pdf",
            ppt_path, "--outdir", task_dir
        ])

        images_from_path = await run_blocking(convert_from_path, pdf_path)
        prs = await run_blocking(PptxPresentation, ppt_path)

        if len(images_from_path) != len(prs.slides):
            raise Exception("PPT页数与生成的图片数量不匹配")

        ppt_video_progress_store[task_id].update(
            total_slides=len(prs.slides), progress_percentage=20.00
        )
        await send_ppt_video_progress(task_id)

        ppt_video_progress_store[task_id].update(current_step=2, progress_percentage=30.00)
        await send_ppt_video_progress(task_id)

        video_segments = []
        with tempfile.TemporaryDirectory() as temp_path:
            for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
                ppt_video_progress_store[task_id].update(
                    current_slide=i + 1,
                    progress_percentage=round(30 + (i / len(prs.slides)) * 40,2)
                )
                await send_ppt_video_progress(task_id)

                notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                if not notes.strip():
                    notes = f"This is the {i + 1} page"

                image_path = pjoin(temp_path, f"frame_{i}.jpg")
                image.save(image_path)

                audio_path = pjoin(temp_path, f"frame_{i}.wav")
                await generate_tts_audio(notes, audio_path)

                video_segment_path = await create_video_segment(
                    image_path, audio_path, temp_path, i
                )
                video_segments.append(video_segment_path)

            ppt_video_progress_store[task_id].update(current_step=3, progress_percentage=80.00)
            await send_ppt_video_progress(task_id)

            output_video_path = pjoin(task_dir, "output.mp4")
            await merge_video_segments(video_segments, output_video_path)

        ppt_video_progress_store[task_id].update(
            status="completed",
            progress_percentage=100.00,
            video_url=f"/api/ppt-to-video/download/{task_id}",
        )
        await send_ppt_video_progress(task_id)

    except Exception as e:
        logger.error(f"PPT2Presentation task failed {task_id}: {e}")
        ppt_video_progress_store[task_id].update(status="failed", error_message=str(e))
        await send_ppt_video_progress(task_id)


async def generate_tts_audio(text: str, output_path: str):
    try:
        sys.path.append(pjoin(os.path.dirname(__file__), "MegaTTS3"))
        from tts.infer_cli import MegaTTS3DiTInfer
        from tts.utils.audio_utils.io import save_wav

        infer = MegaTTS3DiTInfer(ckpt_root=pjoin(os.path.dirname(__file__), "MegaTTS3", "checkpoints", "MegaTTS3"))

        prompt_audio_path = pjoin(os.path.dirname(__file__), "MegaTTS3", "assets", "English_prompt.wav")

        with open(prompt_audio_path, 'rb') as f:
            audio_bytes = f.read()
        latent_file = None
        potential_npy = os.path.splitext(prompt_audio_path)[0] + '.npy'
        if os.path.isfile(potential_npy):
            latent_file = potential_npy
        resource_context = infer.preprocess(audio_bytes, latent_file)

        wav_bytes = infer.forward(
            resource_context,
            text,
            time_step=32,
            p_w=1.6,
            t_w=2.5
        )

        save_wav(wav_bytes, output_path)

    except Exception as e:
        logger.error(f"TTS failed: {str(e)}")
        import numpy as np
        import wave

        sample_rate = 22050
        duration = 3.0
        samples = np.zeros(int(sample_rate * duration), dtype=np.int16)

        with wave.open(output_path, 'w') as wav_file:
            wav_file.setnchannels(1)
            wav_file.setsampwidth(2)
            wav_file.setframerate(sample_rate)
            wav_file.writeframes(samples.tobytes())


async def create_video_segment(image_path: str, audio_path: str, temp_path: str, index: int):
    output_path = pjoin(temp_path, f"segment_{index}.mp4")
    await run_cmd([
        "ffmpeg", "-y", "-loop", "1", "-i", image_path, "-i", audio_path,
        "-vf", "scale=1920:1080", "-c:v", "libx264", "-tune", "stillimage",
        "-c:a", "aac", "-b:a", "192k", "-pix_fmt", "yuv420p", "-shortest",
        output_path
    ])
    return output_path

async def merge_video_segments(video_segments: list[str], output_path: str):
    list_file_path = output_path.replace('.mp4', '_list.txt')
    with open(list_file_path, "w") as f:
        for seg in video_segments:
            f.write(f"file '{seg}'\n")

    await run_cmd([
        "ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", list_file_path,
        "-c", "copy", output_path
    ])
    os.remove(list_file_path)


@app.get("/api/ppt-to-video/download/{task_id}")
async def download_ppt_video(task_id: str):
    if task_id not in ppt_video_progress_store:
        raise HTTPException(status_code=404, detail="Task not found")

    progress = ppt_video_progress_store[task_id]
    if progress["status"] != "completed":
        raise HTTPException(status_code=404, detail="Presentation isn't available")

    video_path = pjoin(progress["task_dir"], "output.mp4")
    if not os.path.exists(video_path):
        raise HTTPException(status_code=404, detail="Presentation not found")

    return FileResponse(
        video_path,
        media_type="video/mp4",
        headers={"Content-Disposition": "attachment; filename=ppt_video.mp4"}
    )


async def ppt_gen(task_id: str, rerun=False):
    if DEBUG:
        importlib.reload(induct)
        importlib.reload(pptgen)
    if rerun:
        task_id = task_id.replace("|", "/")
        active_connections[task_id] = None
        progress_store[task_id] = json.load(open(pjoin(RUNS_DIR, task_id, "task.json")))

    for _ in range(100):
        if task_id in active_connections:
            break
        await asyncio.sleep(0.02)
    else:
        progress_store.pop(task_id)
        return

    task = progress_store.pop(task_id)
    pptx_md5 = task["pptx"]
    pdf_md5 = task["pdf"]
    generation_config = Config(pjoin(RUNS_DIR, task_id))
    pptx_config = Config(pjoin(RUNS_DIR, "pptx", pptx_md5))
    json.dump(task, open(pjoin(generation_config.RUN_DIR, "task.json"), "w"))
    progress = ProgressManager(task_id, STAGES)
    parsedpdf_dir = pjoin(RUNS_DIR, "pdf", pdf_md5)
    ppt_image_folder = pjoin(pptx_config.RUN_DIR, "slide_images")

    await send_progress(
        active_connections[task_id], "task initialized successfully", 10
    )

    try:
        # ppt parsing
        presentation = Presentation.from_file(
            pjoin(pptx_config.RUN_DIR, "source.pptx"), pptx_config
        )
        if not os.path.exists(ppt_image_folder) or len(
                os.listdir(ppt_image_folder)
        ) != len(presentation):
            await ppt_to_images_async(
                pjoin(pptx_config.RUN_DIR, "source.pptx"), ppt_image_folder
            )
            assert len(os.listdir(ppt_image_folder)) == len(presentation) + len(
                presentation.error_history
            ), "Number of parsed slides and images do not match"

            for err_idx, _ in presentation.error_history:
                os.remove(pjoin(ppt_image_folder, f"slide_{err_idx:04d}.jpg"))
            for i, slide in enumerate(presentation.slides, 1):
                slide.slide_idx = i
                os.rename(
                    pjoin(ppt_image_folder, f"slide_{slide.real_idx:04d}.jpg"),
                    pjoin(ppt_image_folder, f"slide_{slide.slide_idx:04d}.jpg"),
                )

        labler = ImageLabler(presentation, pptx_config)
        if os.path.exists(pjoin(pptx_config.RUN_DIR, "image_stats.json")):
            image_stats = json.load(
                open(pjoin(pptx_config.RUN_DIR, "image_stats.json"), encoding="utf-8")
            )
            labler.apply_stats(image_stats)
        else:
            #await labler.caption_images_async(models.vision_model)
            try:
                await labler.caption_images_async(models.vision_model)
            except Exception as e:
                print(f"Skipping image captions due to model error: {e}")
                # [FIX] Force-add captions if they are missing so the app doesn't crash
            for stats in labler.image_stats.values():
                if "caption" not in stats:
                    stats["caption"] = "Image description not available."
            json.dump(
                labler.image_stats,
                open(
                    pjoin(pptx_config.RUN_DIR, "image_stats.json"),
                    "w",
                    encoding="utf-8",
                ),
                ensure_ascii=False,
                indent=4,
            )
        await progress.report_progress()

        # pdf parsing
        if not os.path.exists(pjoin(parsedpdf_dir, "source.md")):
            text_content = parse_pdf(
                pjoin(RUNS_DIR, "pdf", pdf_md5, "source.pdf"),
                parsedpdf_dir,
                models.marker_model,
            )
        else:
            text_content = open(
                pjoin(parsedpdf_dir, "source.md"), encoding="utf-8"
            ).read()
        await progress.report_progress()

        # document refine
        # document refine
        refined_doc_path = pjoin(parsedpdf_dir, "refined_doc.json")
        source_doc = None
        
        if os.path.exists(refined_doc_path):
            try:
                # Try loading to verify integrity
                try:
                    loaded_data = json.load(open(refined_doc_path, encoding="utf-8"))
                except UnicodeDecodeError:
                    loaded_data = json.load(open(refined_doc_path)) # Fallback
                
                source_doc = Document.from_dict(loaded_data, parsedpdf_dir)
            except Exception as e:
                print(f"Warning: Failed to load existing refined_doc.json ({e}). Regenerating...")
                source_doc = None

        if source_doc is None:
            source_doc = await Document.from_markdown_async(
                text_content,
                models.language_model,
                models.vision_model,
                parsedpdf_dir,
            )
            json.dump(
                source_doc.to_dict(),
                open(refined_doc_path, "w", encoding="utf-8"),
                ensure_ascii=False,
                indent=4,
            )
        await progress.report_progress()

        # Slide Induction
        if not os.path.exists(pjoin(pptx_config.RUN_DIR, "slide_induction.json")):
            deepcopy(presentation).save(
                pjoin(pptx_config.RUN_DIR, "template.pptx"), layout_only=True
            )
            await ppt_to_images_async(
                pjoin(pptx_config.RUN_DIR, "template.pptx"),
                pjoin(pptx_config.RUN_DIR, "template_images"),
            )
            slide_inducter = induct.SlideInducterAsync(
                presentation,
                ppt_image_folder,
                pjoin(pptx_config.RUN_DIR, "template_images"),
                pptx_config,
                models.image_model,
                models.language_model,
                models.vision_model,
            )
            layout_induction = await slide_inducter.layout_induct()
            slide_induction = await slide_inducter.content_induct(layout_induction)
            json.dump(
                slide_induction,
                open(
                    pjoin(pptx_config.RUN_DIR, "slide_induction.json"),
                    "w",
                    encoding="utf-8",
                ),
                ensure_ascii=False,
                indent=4,
            )
        else:
            slide_induction = json.load(
                open(
                    pjoin(pptx_config.RUN_DIR, "slide_induction.json"), encoding="utf-8"
                )
            )
        await progress.report_progress()

        # PPT Generation with PPTAgentAsync
        ppt_agent = pptgen.PPTAgentAsync(
            models.text_model,
            models.language_model,
            models.vision_model,
            error_exit=False,
            retry_times=5,
        )
        ppt_agent.set_reference(
            config=generation_config,
            slide_induction=slide_induction,
            presentation=presentation,
        )

        prs, _ = await ppt_agent.generate_pres(
            source_doc=source_doc,
            num_slides=task["numberOfPages"],
        )
        prs.save(pjoin(generation_config.RUN_DIR, "final.pptx"))
        logger.info(f"{task_id}: generation finished")
        await progress.report_progress()
    except Exception as e:
        await progress.fail_stage(str(e))
        traceback.print_exc()


@app.get("/api/evaluate/{task_id}")
async def evaluate_task(task_id: str):
    task_id = task_id.replace("|", "/")
    task_dir = pjoin(RUNS_DIR, task_id)
    if not os.path.exists(task_dir):
        raise HTTPException(status_code=404, detail="Task not found")

    try:
        # Load Task Data
        task = json.load(open(pjoin(task_dir, "task.json")))
        pdf_md5 = task["pdf"]
        pptx_md5 = task["pptx"]
        
        # 1. Get Source Text
        parsedpdf_dir = pjoin(RUNS_DIR, "pdf", pdf_md5)
        source_text_path = pjoin(parsedpdf_dir, "source.md")
        source_text = ""
        if os.path.exists(source_text_path):
             source_text = open(source_text_path, encoding='utf-8').read()

        # 2. Get Generated Content (Outline/Slides) - Simplified to Outline for speed
        slide_induction_path = pjoin(RUNS_DIR, "pptx", pptx_md5, "slide_induction.json")
        presentation_content = ""
        if os.path.exists(slide_induction_path):
            presentation_content = open(slide_induction_path, encoding='utf-8').read()

        # 3. Get Slide Images
        ppt_image_folder = pjoin(RUNS_DIR, "pptx", pptx_md5, "slide_images")
        slide_images = []
        if os.path.exists(ppt_image_folder):
            slide_images = [pjoin(ppt_image_folder, f) for f in sorted(os.listdir(ppt_image_folder)) if f.endswith(".jpg")]

        # Initialize Evaluator (using models from global scope)
        evaluator = Evaluator(models.language_model, models.vision_model)

        # Run Evaluation
        content_score = await evaluator.evaluate_content(source_text, presentation_content)
        visual_score = await evaluator.evaluate_visuals(slide_images)
        quiz = await evaluator.generate_quiz(source_text)

        report = {
            "content_quality": content_score,
            "visual_quality": visual_score,
            "quiz": quiz
        }
        
        # Save Report
        with open(pjoin(task_dir, "evaluation.json"), "w") as f:
            json.dump(report, f, indent=4)
            
        return report

    except Exception as e:
        logger.error(f"Evaluation failed: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn

    ip = "127.0.0.1"
    uvicorn.run(app, host=ip, port=9297)
