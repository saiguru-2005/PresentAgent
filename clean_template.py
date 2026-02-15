from pptx import Presentation
import os

repo_root = os.path.dirname(os.path.abspath(__file__))
template_path = os.path.join(repo_root, "resource", "templates", "default_template.pptx")
output_path = os.path.join(repo_root, "resource", "templates", "clean_template.pptx")

print(f"Cleaning: {template_path}")

if not os.path.exists(template_path):
    print("Template not found!")
    exit(1)

prs = Presentation(template_path)

# Iterate through all slides and shapes to clear text
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            # Clear text but keep the frame
            shape.text_frame.clear()
            # Loop through paragraphs to ensure empty
            for p in shape.text_frame.paragraphs:
                p.text = ""
        if hasattr(shape, "text"):
             if shape.has_text_frame:
                 shape.text_frame.text = ""

# Save as the SAME filename to overwrite the contaminated one (backup first)
backup_path = template_path + ".bak"
if not os.path.exists(backup_path):
    os.rename(template_path, backup_path)
    print(f"Backed up original to {backup_path}")

prs.save(template_path)
print(f"Saved cleaned template to {template_path}")
